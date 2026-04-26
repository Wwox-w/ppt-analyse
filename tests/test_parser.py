"""解析器测试"""

import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from src.parser.pptx_parser import PPTXParser
from src.parser.pdf_parser import ParserFactory


def create_test_pptx(filepath: str):
    """创建一个测试用的 PPTX 文件"""
    prs = Presentation()
    
    # 第一页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "测试课程标题"
    content = slide.placeholders[1]
    content.text = "这是第一页的内容\n- 要点1\n- 要点2"
    
    # 第二页
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "第二章"
    
    prs.save(filepath)


class TestPPTXParser:
    """测试 PPTX 解析器"""

    def setup_method(self):
        self.parser = PPTXParser()
        self.tmp_file = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        create_test_pptx(self.tmp_file.name)
        self.tmp_file.close()

    def teardown_method(self):
        os.unlink(self.tmp_file.name)

    def test_parse_basic(self):
        """测试基本解析功能"""
        result = self.parser.parse(self.tmp_file.name)
        assert result.filename.endswith(".pptx")
        assert result.total_slides == 2
        assert len(result.slides) == 2

    def test_slide_content(self):
        """测试幻灯片内容提取"""
        result = self.parser.parse(self.tmp_file.name)
        first_slide = result.slides[0]
        assert first_slide.slide_number == 1
        assert "测试课程标题" in first_slide.title or "测试课程标题" in first_slide.content

    def test_raw_text(self):
        """测试原始文本提取"""
        result = self.parser.parse(self.tmp_file.name)
        assert "测试课程标题" in result.raw_text
        assert "要点1" in result.raw_text


class TestParserFactory:
    """测试解析器工厂"""

    def test_pptx_parser(self):
        """测试获取 PPTX 解析器"""
        parser = ParserFactory.get_parser("test.pptx")
        from src.parser.pptx_parser import PPTXParser
        assert isinstance(parser, PPTXParser)

    def test_pdf_parser(self):
        """测试获取 PDF 解析器"""
        parser = ParserFactory.get_parser("test.pdf")
        from src.parser.pdf_parser import PDFParser
        assert isinstance(parser, PDFParser)

    def test_invalid_format(self):
        """测试不支持的文件格式"""
        import pytest
        with pytest.raises(ValueError, match="不支持的文件格式"):
            ParserFactory.get_parser("test.doc")
