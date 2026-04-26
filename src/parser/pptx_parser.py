"""PPTX 文件解析器"""

from pptx import Presentation
from pptx.util import Inches, Pt
from src.models.schemas import SlideContent, ParsedPPT


class PPTXParser:
    """解析 .pptx 文件"""

    def parse(self, filepath: str) -> ParsedPPT:
        """
        解析 PPTX 文件

        Args:
            filepath: PPTX 文件路径

        Returns:
            ParsedPPT: 解析后的结构化内容
        """
        prs = Presentation(filepath)
        slides_content = []
        all_text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_data = self._parse_slide(slide, slide_num)
            slides_content.append(slide_data)

            # 收集纯文本用于 LLM 分析
            if slide_data.title:
                all_text_parts.append(f"## 第{slide_num}页: {slide_data.title}")
            else:
                all_text_parts.append(f"## 第{slide_num}页")
            if slide_data.content:
                all_text_parts.append(slide_data.content)
            if slide_data.notes:
                all_text_parts.append(f"[备注]: {slide_data.notes}")

        return ParsedPPT(
            filename=filepath.split("/")[-1],
            total_slides=len(slides_content),
            slides=slides_content,
            raw_text="\n\n".join(all_text_parts),
        )

    def _parse_slide(self, slide, slide_num: int) -> SlideContent:
        """解析单页幻灯片"""
        title = ""
        content_parts = []
        has_tables = False
        has_images = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # 判断是否为标题占位符
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    try:
                        pf = shape.placeholder_format
                        # 通过 idx 判断：标题占位符的 idx 通常为 0
                        if pf.idx == 0:
                            title = text
                            continue
                    except Exception:
                        pass


                # 如果还没有标题，取第一个有意义的文本作为标题
                if not title and len(text) < 100:
                    title = text
                else:
                    content_parts.append(text)

            # 检测表格
            if shape.has_table:
                has_tables = True
                table = shape.table
                table_text = self._extract_table_text(table)
                if table_text:
                    content_parts.append(f"[表格]:\n{table_text}")

            # 检测图片
            if shape.shape_type is not None and "Picture" in str(shape.shape_type):
                has_images = True

        # 提取备注
        notes = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = notes_text

        return SlideContent(
            slide_number=slide_num,
            title=title,
            content="\n".join(content_parts),
            notes=notes,
            has_tables=has_tables,
            has_images=has_images,
        )

    def _extract_table_text(self, table) -> str:
        """提取表格文本为 Markdown 格式"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
