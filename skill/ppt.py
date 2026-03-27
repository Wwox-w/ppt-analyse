from pptx import Presentation

def read_ppt(file_path):
    prs = Presentation(file_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

        slides_content.append({
            "slide": i + 1,
            "content": "\n".join(text)
        })

    return slides_content
import requests

API_URL = "https://api.deepseek.com/v1/chat/completions"
API_KEY = "你的DeepSeekKey"

def analyze_slide(content):
    prompt = f"""
你是一名大学老师，请讲解以下PPT内容：

要求：
1. 说明这一页在讲什么
2. 提取核心知识点
3. 用简单语言解释

PPT内容：
{content}
"""

    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }

    data = {
        "model": "deepseek-chat",
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7
    }

    response = requests.post(API_URL, headers=headers, json=data)
    result = response.json()

    return result["choices"][0]["message"]["content"]
if __name__ == "__main__":
    file_path = "1.pptx"
    result = read_ppt(file_path)

    for slide in result:
        explanation = analyze_slide(slide["content"])

        print(f"Slide {slide['slide']}:")
        print(explanation)
        print("-" * 50)