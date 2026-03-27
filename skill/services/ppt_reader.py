from pptx import Presentation

def read_ppt(file):
    prs = Presentation(file)
    contents = []

    for slide in prs.slides:
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

        contents.append("\n".join(text))

    return contents